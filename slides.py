from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# Create a presentation object
prs = Presentation()

# Title Slide
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
title = slide_title.shapes.title
subtitle = slide_title.placeholders[1]

title.text = "Popular JavaScript Chart Libraries"
subtitle.text = "Overview and Use Cases"

# Data for slides
libraries = [
    ("Chart.js", "chartjs.png", 
     "Features: Simple, easy to use, and highly customizable. Supports a variety of chart types including line, bar, radar, doughnut, and pie charts.\n"
     "Uses: Great for creating simple and interactive charts quickly. Ideal for dashboards, reports, and any application requiring basic data visualization."),
    
    ("D3.js (Data-Driven Documents)", "d3js.png", 
     "Features: Highly flexible and powerful for creating complex and customized visualizations. It uses web standards like SVG, HTML, and CSS.\n"
     "Uses: Best for advanced visualizations and custom charting needs where you need full control over the design and interactivity. Commonly used in data journalism, scientific research, and dynamic data visualizations."),
    
    ("Highcharts", "highcharts.png", 
     "Features: Wide range of chart types and options, supports exporting charts, and has good documentation. Comes with plugins for frameworks like Angular, React, and Vue.\n"
     "Uses: Suitable for creating professional-grade charts and complex data visualizations. Frequently used in business applications, financial reports, and analytics dashboards."),
    
    ("ApexCharts", "apexcharts.png", 
     "Features: Simple to use, good documentation, and supports various chart types like line, bar, heatmap, radar, and more. Integrates well with modern frameworks.\n"
     "Uses: Ideal for modern web applications needing responsive and interactive charts. Commonly used in dashboards and reporting tools."),
    
    ("Plotly.js", "plotlyjs.png", 
     "Features: Supports a wide range of chart types including statistical, scientific, and financial charts. Offers high interactivity and is built on top of D3.js and WebGL.\n"
     "Uses: Excellent for creating highly interactive and detailed charts. Widely used in scientific research, data science projects, and any application requiring advanced plotting capabilities."),
    
    ("ECharts", "echarts.png", 
     "Features: Developed by Baidu, it offers rich visualization options, strong performance, and support for complex data sets. Integrates well with React, Vue, and Angular.\n"
     "Uses: Suitable for applications requiring high-performance and visually appealing charts. Commonly used in business intelligence, data analysis, and interactive dashboards."),
    
    ("Google Charts", "googlecharts.png", 
     "Features: Easy to use, provides a variety of chart types, and integrates well with other Google services. Supports cross-browser compatibility and responsive design.\n"
     "Uses: Great for integrating with Google services and creating simple, interactive charts. Often used in web applications, reports, and presentations."),
    
    ("Recharts", "recharts.png", 
     "Features: Built with React, it offers a simple API, fully componentized charts, and good performance. Easy to integrate with React applications.\n"
     "Uses: Perfect for React developers needing to create charts within their applications. Commonly used in dashboards, data-driven applications, and administrative panels.")
]

# Add a slide for each library with images and improved aesthetics
for lib, img, desc in libraries:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title = slide.shapes.title
    content = slide.placeholders[1]
    left = Inches(5.5)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(f'../adminuser_role_system+{img}', left, top, width=Inches(3), height=Inches(2))
    
    title.text = lib
    content.text = desc

    # Apply color theme and font
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 73, 125)

    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(18)
        paragraph.font.color.rgb = RGBColor(0, 51, 102)

# Save the presentation
pptx_path = "../adminuser_role_system/rr/ptx"
prs.save(pptx_path)

pptx_path
